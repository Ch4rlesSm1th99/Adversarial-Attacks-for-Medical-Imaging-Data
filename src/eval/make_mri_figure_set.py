#!/usr/bin/env python3

from __future__ import annotations
import argparse
import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple

import numpy as np

import visualize as viz

# --------------------------
# utils
# --------------------------

def _ensure_outdir(p: Path) -> Path:
    p.mkdir(parents=True, exist_ok=True)
    return p

def _parse_formats(s: str) -> List[str]:
    return [t.strip().lower() for t in s.split(",") if t.strip()]

def _collect_pngs(outdir: Path) -> List[Path]:
    return sorted(outdir.glob("*.png"))

def _maybe_write_pptx(image_paths: List[Path], pptx_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        print(f"python-pptx not available; skipping PPTX creation: {e}")
        return
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for img in image_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img), Inches(0.5), Inches(0.5), width=Inches(9.0))
    prs.save(str(pptx_path))
    print(f"[OK] Wrote PowerPoint -> {pptx_path}")

def _load_yaml_like_json(p: Path) -> Dict[str, Any]:
    return json.loads(p.read_text())

def _find_subject_cache_dir(processed_dir: Path, names_hint: List[str]) -> Path:
    candidates = [processed_dir] + [d for d in processed_dir.iterdir() if d.is_dir()]
    for d in candidates:
        if (d / f"{names_hint[0]}.npz").exists():
            return d
    # Fallback: look for any dir containing at least one of the names
    for d in candidates:
        if any((d / f"{nm}.npz").exists() for nm in names_hint[:10]):
            return d
    raise SystemExit(f"[ERR] Could not locate subject cache .npz files under {processed_dir}")


class SubjectCacheSliceDataset:
    """
    Iterates 2D slices across subjects stored as per-subject .npz:
      Keys:
        X: shape [4,H,W,D]  or [4,D,H,W]
        seg: shape [H,W,D]  or [D,H,W]
    Returns dict(x=[4,H,W], y=[1], seg=[H,W], subject=str, z=int)
    """

    def __init__(self, cache_root: Path, names: List[str], seed: int = 1337):
        self.cache_root = cache_root
        self.names = names[:]
        self.items: List[Tuple[str, int, int]] = []  # (name, z, label)
        rng = np.random.default_rng(seed)

        # Build an index of all slices with labels from seg
        for name in self.names:
            f = cache_root / f"{name}.npz"
            if not f.exists():
                continue
            with np.load(f) as d:
                seg = d["seg"]
                if seg.ndim != 3:
                    continue
                # Detect axis order for seg
                # If last axis is depth -> seg[..., z]
                # If first axis is depth -> seg[z, ...]
                axes_last = False
                if seg.shape[-1] >= 8 and seg.shape[-1] <= 256:
                    axes_last = True
                # Build pos/neg lists
                D = seg.shape[-1] if axes_last else seg.shape[0]
                pos, neg = [], []
                for z in range(D):
                    s2d = seg[..., z] if axes_last else seg[z]
                    if (s2d > 0).any():
                        pos.append(z)
                    else:
                        neg.append(z)
                # Mix so we do not take all from one subject first
                rng.shuffle(pos); rng.shuffle(neg)
                self.items += [(name, z, 1) for z in pos] + [(name, z, 0) for z in neg]

        # Light shuffle to interleave subjects
        rng.shuffle(self.items)

    def __len__(self) -> int:
        return len(self.items)

    def __getitem__(self, i: int) -> Dict[str, Any]:
        name, z, lab = self.items[i]
        with np.load(self.cache_root / f"{name}.npz") as d:
            X = d["X"]
            seg = d["seg"]
            # Determine axis order
            # Case A: X shape [4,H,W,D], seg [H,W,D] -> z is last axis
            # Case B: X shape [4,D,H,W], seg [D,H,W] -> z is first axis
            if X.shape[-1] == seg.shape[-1] and seg.ndim == 3:
                x = X[..., z]          # [4,H,W]
                s2d = seg[..., z]      # [H,W]
            else:
                x = X[:, z, ...]       # [4,H,W]
                s2d = seg[z, ...]      # [H,W]
            # Normalize to [0,1] per volume dynamic range if needed
            # If data already in [0,1], keep; if looks like 16-bit, scale by 65535
            if np.max(x) > 2.0:
                x = x.astype(np.float32) / 65535.0
            return {
                "x": x.astype(np.float32),
                "y": np.array([float(lab)], dtype=np.float32),
                "seg": s2d.astype(np.uint8),
                "subject": name,
                "z": np.int32(z),
            }

# --------------------------
# multi-format save helpers
# --------------------------

def save_overlay_multi(x: np.ndarray,
                       out_base: Path,
                       formats: Sequence[str],
                       seg: Optional[np.ndarray] = None,
                       ch: int = 0,
                       title: str = "overlay",
                       show_outline: bool = True,
                       outline_th: int = 2,
                       alpha: float = 0.35) -> None:
    for fmt in formats:
        viz.save_overlay_grid(
            x,
            out_path=out_base.with_suffix(f".{fmt}"),
            seg=seg,
            ch=ch,
            title=title,
            show_outline=show_outline,
            outline_th=outline_th,
            alpha=alpha,
        )

def save_modalities_multi(x: np.ndarray,
                          out_base: Path,
                          formats: Sequence[str],
                          ch_names: Sequence[str] = viz.DEFAULT_CH_NAMES,
                          per_row: int = 4,
                          vmin: float = 0.0,
                          vmax: float = 1.0) -> None:
    for fmt in formats:
        viz.save_modalities_grid(
            x,
            out_path=out_base.with_suffix(f".{fmt}"),
            ch_names=ch_names,
            per_row=per_row,
            vmin=vmin,
            vmax=vmax,
        )

def save_fft_multi(x: np.ndarray,
                   out_base: Path,
                   formats: Sequence[str],
                   ch: int = 0,
                   log_scale: bool = True) -> None:
    for fmt in formats:
        viz.save_fft_magnitude_grid(
            x,
            out_path=out_base.with_suffix(f".{fmt}"),
            ch=ch,
            log_scale=log_scale,
        )

# --------------------------
# main
# --------------------------

def main():
    ap = argparse.ArgumentParser(description="Make slide-ready MRI figures from dataset.yaml")
    ap.add_argument("--dataset-yaml", required=True, type=str, help="Path to data/processed/dataset.yaml")
    ap.add_argument("--split", choices=["train", "val", "test"], default="val")
    ap.add_argument("--outdir", required=True, type=str, help="Where to save figures")
    ap.add_argument("--ch", type=int, default=0, help="Channel index (0=FLAIR)")
    ap.add_argument("--num-pos", type=int, default=6, help="How many positive (tumour) samples")
    ap.add_argument("--num-neg", type=int, default=6, help="How many negative (no tumour) samples")
    ap.add_argument("--formats", type=str, default="png,pdf", help="Comma separated: png,pdf")
    ap.add_argument("--pptx", type=str, default="", help="Optional PPTX output path")
    args = ap.parse_args()

    ds_yaml = Path(args.dataset_yaml)
    processed_dir = ds_yaml.parent
    cfg = _load_yaml_like_json(ds_yaml)

    splits_path = processed_dir / "splits.json"
    if not splits_path.exists():
        raise SystemExit(f"[ERR] Missing {splits_path}")

    splits = json.loads(splits_path.read_text())
    names = splits[args.split]
    if not names:
        raise SystemExit(f"[ERR] No subjects found in split '{args.split}' within {splits_path}")

    cache_root = _find_subject_cache_dir(processed_dir, names)
    print(f"[INFO] Using subject cache dir: {cache_root}")

    # Build slice dataset across subjects
    dataset = SubjectCacheSliceDataset(cache_root, names)

    # Pick examples by label
    pos_idx, neg_idx, all_idx = viz.select_examples_by_label(dataset, n_pos=args.num_pos, n_neg=args.num_neg)
    pos_count = min(len(pos_idx), args.num_pos)
    neg_count = min(len(neg_idx), args.num_neg)

    outdir = _ensure_outdir(Path(args.outdir))
    fmts = _parse_formats(args.formats)

    # Tumour overlays with labels
    if pos_count > 0:
        bpos = viz.extract_batch(dataset, pos_idx[:pos_count])
        seg_stack = np.stack(bpos["seg"], 0) if any(s is not None for s in bpos["seg"]) else None
        save_overlay_multi(bpos["x"], outdir / "tumour_overlays_with_labels", fmts,
                           seg=seg_stack, ch=args.ch, title="Tumour overlays", show_outline=True, outline_th=2)

    # Tumour images without labels
    if pos_count > 0:
        bpos = viz.extract_batch(dataset, pos_idx[:pos_count])
        save_overlay_multi(bpos["x"], outdir / "tumour_images_no_labels", fmts,
                           seg=None, ch=args.ch, title="Tumour images (no labels)", show_outline=False)

    # Non-tumour overlays (labels may be absent, but we pass seg if present)
    if neg_count > 0:
        bneg = viz.extract_batch(dataset, neg_idx[:neg_count])
        seg_stack = np.stack(bneg["seg"], 0) if any(s is not None for s in bneg["seg"]) else None
        save_overlay_multi(bneg["x"], outdir / "nontumour_overlays", fmts,
                           seg=seg_stack, ch=args.ch, title="Non-tumour overlays", show_outline=True, outline_th=2)

    # Non-tumour images without labels
    if neg_count > 0:
        bneg = viz.extract_batch(dataset, neg_idx[:neg_count])
        save_overlay_multi(bneg["x"], outdir / "nontumour_images_no_labels", fmts,
                           seg=None, ch=args.ch, title="Non-tumour images (no labels)", show_outline=False)

    # Per-modality grids for a tumour and a non-tumour example
    if pos_count > 0:
        x_pos = dataset[int(pos_idx[0])]["x"]
        save_modalities_multi(x_pos, outdir / "modalities_tumour_example", fmts, ch_names=viz.DEFAULT_CH_NAMES)
    if neg_count > 0:
        x_neg = dataset[int(neg_idx[0])]["x"]
        save_modalities_multi(x_neg, outdir / "modalities_nontumour_example", fmts, ch_names=viz.DEFAULT_CH_NAMES)

    # FFT magnitude for one example (prefer tumour if available)
    x_one = dataset[int(pos_idx[0])]["x"] if pos_count > 0 else dataset[int(neg_idx[0])]["x"]
    save_fft_multi(x_one, outdir / "fft_magnitude_example", fmts, ch=args.ch, log_scale=True)

    # Optional: bundle PNGs into a PowerPoint
    if args.pptx:
        pngs = _collect_pngs(outdir)
        if pngs:
            _maybe_write_pptx(pngs, Path(args.pptx))
        else:
            print("[WARN] No PNGs found to add to PPTX. Skipping.")

    print(f"[DONE] Figures saved under: {outdir}")
    print(f"[INFO] Formats: {', '.join(fmts)}")
    if args.pptx:
        print(f"[INFO] PPTX (if created): {args.pptx}")

if __name__ == "__main__":
    main()
